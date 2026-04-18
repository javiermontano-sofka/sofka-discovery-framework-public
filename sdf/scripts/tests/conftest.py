"""Pytest configuration + fixture builders for the SDF test suite.

Fixtures are synthesised on-demand into `fixtures/` so tests stay reproducible
and the repository doesn't need to ship binary blobs (PDF, DOCX, PPTX).
"""
from __future__ import annotations

import os
import shutil
import subprocess
import sys
from pathlib import Path

import pytest

# Make the ecosystem helpers importable from tests
HERE = Path(__file__).resolve().parent
PLUGIN_ROOT = HERE.parent.parent
SCRIPTS = PLUGIN_ROOT / "scripts"
sys.path.insert(0, str(SCRIPTS))
sys.path.insert(0, str(SCRIPTS / "ecosystem"))


# Ensure venv exists before tests
@pytest.fixture(scope="session", autouse=True)
def ensure_venv():
    venv = SCRIPTS / ".venv"
    if not venv.is_dir():
        subprocess.run(["bash", str(SCRIPTS / "setup-attachments.sh")], check=True)
    yield


def _venv_python() -> str:
    return str(SCRIPTS / ".venv" / "bin" / "python")


@pytest.fixture
def plugin_root() -> Path:
    return PLUGIN_ROOT


@pytest.fixture
def venv_python() -> str:
    return _venv_python()


@pytest.fixture
def fixtures_dir() -> Path:
    return HERE / "fixtures"


# --- Synthetic fixture builders ---------------------------------------------

@pytest.fixture
def csv_path(tmp_path: Path) -> Path:
    p = tmp_path / "smoke.csv"
    p.write_text(
        "id,country,amount\n"
        "1,MX,100\n2,CO,200\n3,PE,150\n4,CL,275\n5,AR,90\n",
        encoding="utf-8",
    )
    return p


@pytest.fixture
def xlsx_path(tmp_path: Path) -> Path:
    from openpyxl import Workbook
    p = tmp_path / "smoke.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.title = "Budget"
    ws.append(["Phase", "P50", "P80"])
    ws.append(["Discover", 12, 18])
    ws.append(["Design", 30, 42])
    ws2 = wb.create_sheet("Stakeholders")
    ws2.append(["Role", "Country"])
    ws2.append(["CIO", "MX"])
    wb.save(p)
    return p


@pytest.fixture
def docx_path(tmp_path: Path) -> Path:
    from docx import Document
    p = tmp_path / "smoke.docx"
    d = Document()
    d.add_heading("Discovery Brief", 1)
    d.add_paragraph("Cliente: AcmeCorp.")
    d.add_heading("Scope", 2)
    d.add_paragraph("Retail LatAm, 200 usuarios.")
    table = d.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "Module"
    table.rows[0].cells[1].text = "Status"
    table.rows[1].cells[0].text = "FI"
    table.rows[1].cells[1].text = "Active"
    d.save(p)
    return p


@pytest.fixture
def pdf_path(tmp_path: Path) -> Path:
    from pypdf import PdfWriter
    from pypdf.generic import RectangleObject
    p = tmp_path / "smoke.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=595, height=842)
    writer.add_blank_page(width=595, height=842)
    with open(p, "wb") as f:
        writer.write(f)
    return p


@pytest.fixture
def pptx_path(tmp_path: Path) -> Path:
    from pptx import Presentation
    p = tmp_path / "smoke.pptx"
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    s1 = prs.slides.add_slide(slide_layout)
    s1.shapes.title.text = "Kickoff"
    s2 = prs.slides.add_slide(slide_layout)
    s2.shapes.title.text = "Scope"
    if s2.has_notes_slide:
        s2.notes_slide.notes_text_frame.text = "Note: validar con Finance"
    prs.save(p)
    return p


@pytest.fixture
def html_path(tmp_path: Path) -> Path:
    p = tmp_path / "smoke.html"
    p.write_text(
        "<html><head><title>Demo</title></head><body>"
        "<h1>Hello</h1><p>World</p>"
        "<table><tr><th>k</th><th>v</th></tr><tr><td>a</td><td>1</td></tr></table>"
        "<a href='https://example.com'>example</a>"
        "</body></html>",
        encoding="utf-8",
    )
    return p


@pytest.fixture
def py_path(tmp_path: Path) -> Path:
    p = tmp_path / "smoke.py"
    p.write_text(
        "import os\n"
        "from pathlib import Path\n\n"
        "class Foo:\n    def bar(self, x):\n        return x * 2\n\n"
        "def main():\n    print('hi')\n",
        encoding="utf-8",
    )
    return p


@pytest.fixture
def sql_path(tmp_path: Path) -> Path:
    p = tmp_path / "smoke.sql"
    p.write_text(
        "CREATE TABLE customers (id INT, name VARCHAR(100));\n"
        "ALTER TABLE customers ADD COLUMN country CHAR(2);\n",
        encoding="utf-8",
    )
    return p


@pytest.fixture
def json_path(tmp_path: Path) -> Path:
    import json
    p = tmp_path / "smoke.json"
    p.write_text(json.dumps({"client": "Acme", "modules": ["FI", "CO"]}), encoding="utf-8")
    return p


@pytest.fixture
def yaml_path(tmp_path: Path) -> Path:
    p = tmp_path / "smoke.yaml"
    p.write_text("client: Acme\nmodules:\n  - FI\n  - CO\n", encoding="utf-8")
    return p


@pytest.fixture
def md_path(tmp_path: Path) -> Path:
    p = tmp_path / "demo.md"
    p.write_text(
        "# Comité SAP — Demo\n\n"
        "> TL;DR: prueba render\n\n"
        "## Hallazgos\n\n"
        "- Custom code debt alto [CÓDIGO]\n"
        "- Datos incompletos [ADJUNTO:customers.csv:col=Country]\n"
        "- Recomendación Bluefield [INFERENCIA]\n\n"
        "## Tabla\n\n"
        "| Fase | P50 | P80 |\n|------|-----|-----|\n| Prepare | 24 | 32 |\n",
        encoding="utf-8",
    )
    return p


# --- Generic helper to invoke an extractor through the venv -----------------

def run_extractor(extractor: str, input_path: Path, out_path: Path | None = None,
                  cwd: Path | None = None) -> dict:
    """Run scripts/extract-<extractor>.py via venv python and return parsed metadata.

    Returns dict with keys: returncode, stdout, stderr, output_path, output_text.
    """
    script = SCRIPTS / f"extract-{extractor}.py"
    out = out_path or input_path.with_suffix(".priming.md")
    env = os.environ.copy()
    env["PYTHONPATH"] = f"{SCRIPTS}:{env.get('PYTHONPATH','')}"
    res = subprocess.run(
        [_venv_python(), str(script), str(input_path), "--out", str(out)],
        cwd=str(cwd or input_path.parent),
        env=env,
        capture_output=True,
        text=True,
    )
    return {
        "returncode": res.returncode,
        "stdout": res.stdout,
        "stderr": res.stderr,
        "output_path": out,
        "output_text": out.read_text(encoding="utf-8") if out.is_file() else "",
    }
