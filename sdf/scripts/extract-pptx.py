#!/usr/bin/env python3
"""PPTX → priming-rag. Slides + speaker notes."""
from pathlib import Path
from pptx import Presentation
from _common import cli, file_meta, write_priming, truncate


def main() -> None:
    args = cli()
    path = Path(args.path)
    meta = file_meta(path)
    prs = Presentation(str(path))

    blocks: list[str] = []
    titles: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        title = ""
        body: list[str] = []
        for shp in slide.shapes:
            if not shp.has_text_frame:
                continue
            for p in shp.text_frame.paragraphs:
                text = "".join(r.text for r in p.runs).strip()
                if not text:
                    continue
                if not title and (shp.name or "").lower().startswith("title"):
                    title = text
                else:
                    body.append(text)
        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        titles.append(title or f"Slide {i}")
        block = f"### Slide {i}: {title}\n"
        if body:
            block += "\n".join(f"- {b}" for b in body)
        if notes:
            block += f"\n\n> **Notas:** {notes}"
        blocks.append(block)

    resumen = [
        f"Deck con {len(prs.slides)} slides",
        "Títulos: " + "; ".join(titles[:6]),
    ]
    contenido = "\n\n".join(blocks)
    evidencia = [f"[ADJUNTO:{path.name}:slide={i}]" for i in range(1, min(len(prs.slides), 6) + 1)]
    out = write_priming(Path(args.out) if args.out else None, meta, "pptx", resumen, truncate(contenido), evidencia)
    print(str(out))


if __name__ == "__main__":
    main()
